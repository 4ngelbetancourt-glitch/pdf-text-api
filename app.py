import io
import os
import re
import tempfile
import time
import zipfile
import logging
from typing import Any, Dict, List, Optional, Tuple

import filetype
import py7zr
import rarfile
from fastapi import FastAPI, File, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse

# ----------------------------------------------------------------------
# Logging configuration
# ----------------------------------------------------------------------
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("text_extractor")

# ----------------------------------------------------------------------
# FastAPI application setup
# ----------------------------------------------------------------------
app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ----------------------------------------------------------------------
# Constants
# ----------------------------------------------------------------------
MAX_UPLOAD_SIZE = 200 * 1024 * 1024  # 200 MB
MAX_EXTRACTED_FILE_SIZE = 200 * 1024 * 1024  # 200 MB per extracted file
MAX_ARCHIVE_DEPTH = 3
MAX_TOTAL_FILES = 2000

# ----------------------------------------------------------------------
# Configure unrar binary path
# ----------------------------------------------------------------------
UNRAR_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "unrar")
if os.path.exists(UNRAR_PATH):
    rarfile.UNRAR_TOOL = UNRAR_PATH
    logger.info(f"Usando binario unrar en: {UNRAR_PATH}")
else:
    logger.warning("Binario unrar no encontrado, el soporte RAR dependerá del PATH del sistema.")

# ----------------------------------------------------------------------
# Extraction context to track global limits across recursive calls
# ----------------------------------------------------------------------
class ExtractionContext:
    def __init__(self, max_depth: int = MAX_ARCHIVE_DEPTH,
                 max_files: int = MAX_TOTAL_FILES,
                 max_size: int = MAX_UPLOAD_SIZE):
        self.max_depth = max_depth
        self.max_files = max_files
        self.max_size = max_size
        self.files_processed: int = 0
        self.files_failed: int = 0

# ----------------------------------------------------------------------
# Text cleaning
# ----------------------------------------------------------------------
def clean_text(text: str) -> str:
    """Normalize text while preserving meaningful structure."""
    if not text:
        return ""
    # Remove null bytes
    text = text.replace("\x00", "")
    # Normalize line endings
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    # Collapse multiple spaces (but not newlines)
    text = re.sub(r"[ \t]+", " ", text)
    # Remove excessive blank lines (more than 2 consecutive)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()

# ----------------------------------------------------------------------
# Format detection (content first, then magic bytes, then extension, then MIME)
# ----------------------------------------------------------------------
def detect_format(file_path: str, original_filename: str) -> str:
    """
    Return a canonical format identifier (lowercase) like 'pdf', 'docx', 'zip'.
    Uses content sniffing (filetype), manual magic bytes, fallback to extension, then MIME type.
    """
    # 1. Content detection
    kind = filetype.guess(file_path)
    if kind is not None:
        mime = kind.mime
        ext = kind.extension
        # Map MIME/extension to our identifiers
        mime_map = {
            "application/pdf": "pdf",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
            "application/vnd.oasis.opendocument.text": "odt",
            "application/vnd.oasis.opendocument.spreadsheet": "ods",
            "application/zip": "zip",
            "application/x-7z-compressed": "7z",
            "application/x-rar-compressed": "rar",
            "application/vnd.rar": "rar",
            "application/rtf": "rtf",
            "text/plain": "txt",
            "text/html": "html",
            "text/csv": "csv",
            "application/json": "json",
            "application/xml": "xml",
            "text/xml": "xml",
            "text/markdown": "md",
        }
        # Try by MIME first
        if mime in mime_map:
            return mime_map[mime]
        # Try by extension returned by filetype
        if ext in ("pdf", "docx", "xlsx", "pptx", "odt", "ods",
                    "zip", "7z", "rar", "rtf", "txt", "html", "csv", "json",
                    "xml", "md", "log"):
            return ext
        # Try extension from original filename (before falling through)
        ext = os.path.splitext(original_filename)[1].lower().lstrip(".")
        if ext in ("pdf", "docx", "xlsx", "pptx", "odt", "ods",
                    "zip", "7z", "rar", "rtf", "txt", "html", "csv", "json",
                    "xml", "md", "log"):
            return ext
        # Do NOT return "unknown" here – let it fall through to magic bytes and MIME fallback

    # 2. Manual magic bytes check for formats not covered by filetype (e.g., RAR)
    try:
        with open(file_path, 'rb') as f:
            header = f.read(4)
        if header[:4] == b'Rar!':
            return "rar"
    except Exception:
        pass

    # 3. Fallback to extension from original filename (already checked above, but just in case)
    ext = os.path.splitext(original_filename)[1].lower().lstrip(".")
    ext_map = {
        "pdf": "pdf",
        "docx": "docx",
        "xlsx": "xlsx",
        "pptx": "pptx",
        "odt": "odt",
        "ods": "ods",
        "zip": "zip",
        "7z": "7z",
        "rar": "rar",
        "rtf": "rtf",
        "txt": "txt",
        "html": "html",
        "htm": "html",
        "csv": "csv",
        "json": "json",
        "xml": "xml",
        "md": "md",
        "log": "log",
    }
    if ext in ext_map:
        return ext_map[ext]

    # 4. Last resort: MIME from filename
    import mimetypes
    mime, _ = mimetypes.guess_type(original_filename)
    if mime:
        if "pdf" in mime: return "pdf"
        if "word" in mime or "docx" in mime: return "docx"
        if "spreadsheet" in mime or "xlsx" in mime: return "xlsx"
        if "presentation" in mime or "pptx" in mime: return "pptx"
        if "opendocument.text" in mime: return "odt"
        if "opendocument.spreadsheet" in mime: return "ods"
        if "zip" in mime: return "zip"
        if "x-7z" in mime: return "7z"
        if "rar" in mime: return "rar"
        if "rtf" in mime: return "rtf"
        if "plain" in mime: return "txt"
        if "html" in mime: return "html"
        if "csv" in mime: return "csv"
        if "json" in mime: return "json"
        if "xml" in mime: return "xml"
        if "markdown" in mime: return "md"
    return "unknown"

# ----------------------------------------------------------------------
# Individual extractors
# ----------------------------------------------------------------------
def extract_pdf(file_path: str) -> Tuple[str, dict]:
    """Extract text from PDF, trying empty password if encrypted."""
    from pypdf import PdfReader
    reader = PdfReader(file_path)
    if reader.is_encrypted:
        try:
            reader.decrypt("")
        except Exception:
            raise ValueError("PDF is encrypted and cannot be decrypted with an empty password.")
    text_parts = []
    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text_parts.append(page_text)
    text = "\n".join(text_parts)
    metadata = {"pages": len(reader.pages)}
    return text, metadata

def extract_docx(file_path: str) -> Tuple[str, dict]:
    """Extract text and table content from DOCX."""
    from docx import Document
    doc = Document(file_path)
    parts = []
    for para in doc.paragraphs:
        parts.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            row_text = "\t".join(cell.text for cell in row.cells)
            parts.append(row_text)
    text = "\n".join(parts)
    return text, {}

def extract_pptx(file_path: str) -> Tuple[str, dict]:
    """Extract text from all slides of a PPTX."""
    from pptx import Presentation
    prs = Presentation(file_path)
    slides_text = []
    for slide in prs.slides:
        slide_parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    slide_parts.append(paragraph.text)
        if slide_parts:
            slides_text.append("\n".join(slide_parts))
    text = "\n\n".join(slides_text)
    return text, {"slides": len(prs.slides)}

def extract_xlsx(file_path: str) -> Tuple[str, dict]:
    """Extract text from all sheets of an XLSX."""
    from openpyxl import load_workbook
    wb = load_workbook(file_path, read_only=True, data_only=True)
    sheet_texts = []
    for name in wb.sheetnames:
        ws = wb[name]
        rows_text = []
        for row in ws.iter_rows(values_only=True):
            # Convert all cells to strings, skip fully empty rows
            row_values = [str(cell) if cell is not None else "" for cell in row]
            if any(v.strip() for v in row_values):
                rows_text.append("\t".join(row_values))
        if rows_text:
            sheet_texts.append(f"--- Sheet: {name} ---\n" + "\n".join(rows_text))
    wb.close()
    text = "\n\n".join(sheet_texts)
    return text, {"sheets": len(wb.sheetnames)}

def extract_html(file_path: str) -> Tuple[str, dict]:
    """Extract visible text from HTML."""
    from bs4 import BeautifulSoup
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        soup = BeautifulSoup(f.read(), "html.parser")
    # Remove script and style elements
    for tag in soup(["script", "style"]):
        tag.decompose()
    text = soup.get_text(separator="\n")
    return text, {}

def extract_json(file_path: str) -> Tuple[str, dict]:
    """Serialize JSON to an indented string."""
    import json
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        data = json.load(f)
    text = json.dumps(data, indent=2, ensure_ascii=False)
    return text, {}

def extract_xml(file_path: str) -> Tuple[str, dict]:
    """Extract all text content from XML."""
    import xml.etree.ElementTree as ET
    tree = ET.parse(file_path)
    root = tree.getroot()
    # Collect all text recursively
    parts = []
    for elem in root.iter():
        if elem.text:
            parts.append(elem.text)
        if elem.tail:
            parts.append(elem.tail)
    text = " ".join(parts)
    return text, {}

def extract_rtf(file_path: str) -> Tuple[str, dict]:
    """Convert RTF to plain text."""
    from striprtf.striprtf import rtf_to_text
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        rtf_content = f.read()
    text = rtf_to_text(rtf_content)
    return text, {}

def extract_odt(file_path: str) -> Tuple[str, dict]:
    """Extract text from ODT by reading content.xml inside the ZIP."""
    import zipfile
    import xml.etree.ElementTree as ET
    with zipfile.ZipFile(file_path, 'r') as z:
        with z.open('content.xml') as xml_file:
            tree = ET.parse(xml_file)
    root = tree.getroot()
    # Namespace handling
    ns = {'text': 'urn:oasis:names:tc:opendocument:xmlns:text:1.0',
          'office': 'urn:oasis:names:tc:opendocument:xmlns:office:1.0'}
    paragraphs = root.findall('.//text:p', ns) + root.findall('.//text:h', ns)
    parts = []
    for p in paragraphs:
        # Extract text from element and its children
        text = "".join(p.itertext())
        parts.append(text)
    return "\n".join(parts), {}

def extract_ods(file_path: str) -> Tuple[str, dict]:
    """Extract tabular content from ODS via content.xml."""
    import zipfile
    import xml.etree.ElementTree as ET
    with zipfile.ZipFile(file_path, 'r') as z:
        with z.open('content.xml') as xml_file:
            tree = ET.parse(xml_file)
    root = tree.getroot()
    ns = {
        'table': 'urn:oasis:names:tc:opendocument:xmlns:table:1.0',
        'text': 'urn:oasis:names:tc:opendocument:xmlns:text:1.0',
        'office': 'urn:oasis:names:tc:opendocument:xmlns:office:1.0'
    }
    tables = root.findall('.//table:table', ns)
    all_table_texts = []
    for table in tables:
        rows = table.findall('.//table:table-row', ns)
        table_rows = []
        for row in rows:
            cells = row.findall('.//table:table-cell', ns)
            cell_texts = []
            for cell in cells:
                # Get text from all text:p inside cell
                text_ps = cell.findall('.//text:p', ns)
                cell_content = " ".join("".join(p.itertext()) for p in text_ps)
                cell_texts.append(cell_content)
            table_rows.append("\t".join(cell_texts))
        if table_rows:
            all_table_texts.append("\n".join(table_rows))
    return "\n\n".join(all_table_texts), {}

def extract_txt(file_path: str) -> Tuple[str, dict]:
    """Read plain text file with encoding detection fallback."""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            text = f.read()
    except UnicodeDecodeError:
        with open(file_path, "r", encoding="latin-1") as f:
            text = f.read()
    return text, {}

# Extractor registry
EXTRACTORS: Dict[str, Any] = {
    "pdf": extract_pdf,
    "docx": extract_docx,
    "pptx": extract_pptx,
    "xlsx": extract_xlsx,
    "html": extract_html,
    "json": extract_json,
    "xml": extract_xml,
    "rtf": extract_rtf,
    "odt": extract_odt,
    "ods": extract_ods,
    "txt": extract_txt,
    "md": extract_txt,   # same as plain text
    "log": extract_txt,
    "csv": extract_txt,  # CSV is plain text; optionally we could parse with csv module, but plain text works
}

# ----------------------------------------------------------------------
# Archive processing (ZIP, 7z, RAR)
# ----------------------------------------------------------------------
def process_archive(file_path: str, archive_type: str,
                    original_filename: str, depth: int,
                    context: ExtractionContext) -> Dict[str, Any]:
    """
    Extract archive, recursively process all contained files,
    and return combined text with headers.
    """
    combined_text_parts = []
    processed = 0
    failed = 0
    temp_dir = tempfile.mkdtemp(prefix="archive_")

    try:
        if archive_type == "zip":
            with zipfile.ZipFile(file_path, 'r') as zf:
                members = [m for m in zf.infolist() if not m.is_dir()]
                for member in members:
                    # Check global limit
                    if context.files_processed + context.files_failed >= context.max_files:
                        logger.warning("Reached maximum file limit, skipping remaining entries.")
                        break
                    if depth > context.max_depth:
                        logger.warning(f"Max depth exceeded for {member.filename}, skipping.")
                        failed += 1
                        continue
                    if member.file_size > context.max_size:
                        logger.warning(f"File too large inside archive: {member.filename}")
                        failed += 1
                        continue
                    # Extract single file
                    extracted_path = os.path.join(temp_dir, member.filename)
                    os.makedirs(os.path.dirname(extracted_path), exist_ok=True)
                    with zf.open(member) as source, open(extracted_path, 'wb') as target:
                        # Read in chunks to avoid memory blow
                        while True:
                            chunk = source.read(1024 * 1024)
                            if not chunk:
                                break
                            target.write(chunk)
                    # Process extracted file
                    result = process_file(
                        file_path=extracted_path,
                        original_filename=member.filename,
                        depth=depth + 1,
                        context=context
                    )
                    if result["success"]:
                        header = f"==================================================\nARCHIVO: {member.filename}\n=================================================="
                        combined_text_parts.append(f"{header}\n\n{result['text']}")
                        processed += 1
                        context.files_processed += 1
                    else:
                        failed += 1
                        context.files_failed += 1
                    # Clean up extracted file immediately
                    try:
                        os.remove(extracted_path)
                    except Exception:
                        pass

        elif archive_type == "7z":
            with py7zr.SevenZipFile(file_path, mode='r') as szf:
                members = szf.getnames()
                for member_name in members:
                    # Check if member is a file (not directory)
                    info = szf.getinfo(member_name)
                    if info.is_directory():
                        continue
                    if context.files_processed + context.files_failed >= context.max_files:
                        break
                    if depth > context.max_depth:
                        failed += 1
                        continue
                    if info.size > context.max_size:
                        failed += 1
                        continue
                    # Extract single file to temp
                    extracted_path = os.path.join(temp_dir, member_name)
                    os.makedirs(os.path.dirname(extracted_path), exist_ok=True)
                    # Extract only this file
                    szf.extract(targets=[member_name], path=temp_dir)
                    # Process
                    result = process_file(
                        file_path=extracted_path,
                        original_filename=member_name,
                        depth=depth + 1,
                        context=context
                    )
                    if result["success"]:
                        header = f"==================================================\nARCHIVO: {member_name}\n=================================================="
                        combined_text_parts.append(f"{header}\n\n{result['text']}")
                        processed += 1
                        context.files_processed += 1
                    else:
                        failed += 1
                        context.files_failed += 1
                    # Cleanup
                    try:
                        os.remove(extracted_path)
                    except Exception:
                        pass

        elif archive_type == "rar":
            with rarfile.RarFile(file_path, 'r') as rf:
                members = [m for m in rf.infolist() if not m.is_dir()]
                for member in members:
                    # Check global limit
                    if context.files_processed + context.files_failed >= context.max_files:
                        logger.warning("Reached maximum file limit, skipping remaining entries in RAR.")
                        break
                    if depth > context.max_depth:
                        logger.warning(f"Max depth exceeded for {member.filename}, skipping.")
                        failed += 1
                        continue
                    if member.file_size > context.max_size:
                        logger.warning(f"File too large inside RAR: {member.filename}")
                        failed += 1
                        continue
                    # Extract single file
                    extracted_path = os.path.join(temp_dir, member.filename)
                    os.makedirs(os.path.dirname(extracted_path), exist_ok=True)
                    try:
                        rf.extract(member, path=os.path.dirname(extracted_path))
                    except rarfile.RarWrongPassword:
                        logger.warning(f"Wrong password for RAR entry: {member.filename}, skipping.")
                        failed += 1
                        context.files_failed += 1
                        continue
                    except rarfile.RarCryptoError:
                        logger.warning(f"Encryption error in RAR entry: {member.filename}, skipping.")
                        failed += 1
                        context.files_failed += 1
                        continue
                    # Process extracted file
                    result = process_file(
                        file_path=extracted_path,
                        original_filename=member.filename,
                        depth=depth + 1,
                        context=context
                    )
                    if result["success"]:
                        header = f"==================================================\nARCHIVO: {member.filename}\n=================================================="
                        combined_text_parts.append(f"{header}\n\n{result['text']}")
                        processed += 1
                        context.files_processed += 1
                    else:
                        failed += 1
                        context.files_failed += 1
                    # Clean up extracted file immediately
                    try:
                        os.remove(extracted_path)
                    except Exception:
                        pass

        else:
            raise ValueError(f"Unsupported archive type: {archive_type}")

    except Exception as e:
        logger.error(f"Error processing archive {original_filename}: {e}")
        return {"text": "", "type": archive_type, "success": False,
                "metadata": {}, "processed": 0, "failed": 1}
    finally:
        # Remove entire temp directory
        try:
            import shutil
            shutil.rmtree(temp_dir, ignore_errors=True)
        except Exception:
            pass

    combined_text = "\n\n".join(combined_text_parts)
    return {
        "text": combined_text,
        "type": archive_type,
        "success": True,
        "metadata": {},
        "processed": processed,
        "failed": failed
    }

# ----------------------------------------------------------------------
# Universal file processor (dispatcher)
# ----------------------------------------------------------------------
def process_file(file_path: str, original_filename: str,
                 depth: int, context: ExtractionContext) -> Dict[str, Any]:
    """
    Detect format, extract text, handle archives recursively.
    Returns a dict with text, type, success, metadata, and optional counts
    (for single files, processed/failed are 1/0).
    """
    # 0. Check if we already reached global limits
    if context.files_processed + context.files_failed >= context.max_files:
        return {"text": "", "type": "unknown", "success": False,
                "metadata": {}}

    # Detect format
    fmt = detect_format(file_path, original_filename)
    logger.info(f"Processing file: {original_filename} (detected as {fmt})")

    # Handle archives recursively
    if fmt in ("zip", "7z", "rar"):
        result = process_archive(file_path, fmt, original_filename, depth, context)
        # Update context counters (archive function already updated context.files_processed/failed)
        return result

    # If format has a known extractor
    if fmt in EXTRACTORS:
        try:
            text, metadata = EXTRACTORS[fmt](file_path)
            text = clean_text(text)
            context.files_processed += 1
            return {
                "text": text,
                "type": fmt,
                "success": True,
                "metadata": metadata
            }
        except Exception as e:
            logger.error(f"Failed to extract {original_filename}: {e}")
            context.files_failed += 1
            return {"text": "", "type": fmt, "success": False, "metadata": {}}

    # Unknown format: try generic text reading
    try:
        with open(file_path, "rb") as f:
            raw = f.read()
        # Try UTF-8, then latin-1
        try:
            text = raw.decode("utf-8")
        except UnicodeDecodeError:
            text = raw.decode("latin-1")
        text = clean_text(text)
        if text:
            context.files_processed += 1
            return {"text": text, "type": "txt", "success": True, "metadata": {}}
        else:
            raise ValueError("Empty content")
    except Exception as e:
        logger.error(f"Could not read as text: {original_filename} - {e}")
        context.files_failed += 1
        return {"text": "", "type": "unknown", "success": False, "metadata": {}}

# ----------------------------------------------------------------------
# API endpoints
# ----------------------------------------------------------------------
@app.get("/")
def root():
    return {"status": "ok"}

@app.post("/extract-text")
async def extract_text(file: UploadFile = File(...)):
    start_time = time.time()
    temp_path = None
    try:
        # 1. Stream uploaded file to temporary file
        size = 0
        suffix = os.path.splitext(file.filename or "file.bin")[1] or ".tmp"
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
            temp_path = tmp.name
            while True:
                chunk = await file.read(1024 * 1024)  # 1 MB
                if not chunk:
                    break
                size += len(chunk)
                if size > MAX_UPLOAD_SIZE:
                    tmp.close()
                    os.remove(temp_path)
                    return JSONResponse(
                        {"error": "File exceeds maximum size of 200 MB"},
                        status_code=413,
                    )
                tmp.write(chunk)

        # 2. Process file with extraction context
        context = ExtractionContext()
        result = process_file(
            file_path=temp_path,
            original_filename=file.filename or "unnamed",
            depth=1,
            context=context,
        )

        # 3. Build response
        processing_time = time.time() - start_time
        response = {
            "filename": file.filename,
            "type": result["type"],
            "size_mb": round(size / (1024 * 1024), 2),
            "characters": len(result["text"]),
            "files_processed": context.files_processed,
            "files_failed": context.files_failed,
            "processing_time_seconds": round(processing_time, 3),
            "text": result["text"],
        }
        # Include PDF pages if available (backward compatibility)
        if result["type"] == "pdf" and "pages" in result.get("metadata", {}):
            response["pages"] = result["metadata"]["pages"]

        return response

    except Exception as e:
        logger.exception("Unhandled error during extraction")
        return JSONResponse(
            {"error": str(e)},
            status_code=500,
        )

    finally:
        if temp_path and os.path.exists(temp_path):
            try:
                os.remove(temp_path)
            except Exception:
                pass